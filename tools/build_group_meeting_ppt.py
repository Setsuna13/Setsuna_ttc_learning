from pathlib import Path
import zipfile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


OUT = Path("output")
OUT.mkdir(exist_ok=True)
FIG_DIR = OUT / "assets" / "figures"
FIG_DIR.mkdir(parents=True, exist_ok=True)
PPTX_PATH = OUT / "final_presentation_cn.pptx"
QA_PATH = OUT / "qa_report.md"
NOTES_PATH = OUT / "speaker_notes_cn.md"
OUTLINE_PATH = OUT / "ppt_outline_cn.md"


W, H = Inches(13.333), Inches(7.5)


COLORS = {
    "ink": RGBColor(30, 37, 48),
    "muted": RGBColor(92, 102, 115),
    "light": RGBColor(246, 248, 250),
    "line": RGBColor(205, 213, 224),
    "blue": RGBColor(35, 100, 170),
    "teal": RGBColor(0, 132, 120),
    "orange": RGBColor(204, 108, 48),
    "red": RGBColor(184, 68, 68),
    "green": RGBColor(48, 132, 82),
    "purple": RGBColor(112, 82, 160),
    "yellow": RGBColor(248, 226, 142),
    "white": RGBColor(255, 255, 255),
}


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color=COLORS["line"], width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_textbox(slide, x, y, w, h, text, size=16, bold=False, color=None,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="Noto Sans CJK SC"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = font
    p.font.color.rgb = color or COLORS["ink"]
    return box


def add_multiline(slide, x, y, w, h, lines, size=14, color=None, bullet=False,
                  gap=0.8, font="Noto Sans CJK SC"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(5 * gap)
        p.font.size = Pt(size)
        p.font.name = font
        p.font.color.rgb = color or COLORS["ink"]
        if bullet:
            p.text = "· " + line
    return box


def add_title(slide, title, kicker=None):
    if kicker:
        add_textbox(slide, Inches(0.55), Inches(0.22), Inches(4.0), Inches(0.28),
                    kicker, size=8.5, bold=True, color=COLORS["blue"])
    add_textbox(slide, Inches(0.55), Inches(0.42), Inches(11.9), Inches(0.55),
                title, size=23, bold=True, color=COLORS["ink"])
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.03), Inches(12.2), Inches(0.015))
    set_fill(line, COLORS["line"])
    line.line.fill.background()


def add_source(slide, text):
    add_textbox(slide, Inches(0.55), Inches(7.10), Inches(12.2), Inches(0.2),
                text, size=7.5, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_tag(slide, x, y, text, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.55), Inches(0.32))
    set_fill(shape, color)
    shape.line.fill.background()
    add_textbox(slide, x + Inches(0.06), y + Inches(0.04), Inches(1.43), Inches(0.22),
                text, size=8.5, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
    return shape


def add_card(slide, x, y, w, h, title, body, color=COLORS["blue"], fill=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(shape, fill or COLORS["white"])
    set_line(shape, COLORS["line"], 0.8)
    add_textbox(slide, x + Inches(0.18), y + Inches(0.13), w - Inches(0.36), Inches(0.25),
                title, size=12, bold=True, color=color)
    add_multiline(slide, x + Inches(0.18), y + Inches(0.50), w - Inches(0.36), h - Inches(0.65),
                  body if isinstance(body, list) else [body], size=10.8, color=COLORS["ink"])
    return shape


def arrow(slide, x1, y1, x2, y2, color=COLORS["muted"], width=1.6):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.line.end_arrowhead = True
    return conn


def make_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]
    slides = []

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    slides.append("封面")
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    set_fill(bg, COLORS["light"])
    bg.line.fill.background()
    add_tag(slide, Inches(0.62), Inches(0.62), "组会汇报", COLORS["blue"])
    add_textbox(slide, Inches(0.62), Inches(1.35), Inches(11.4), Inches(1.35),
                "从 TTC 到碰撞概率：\nPNR 与 Affordance-Attention 的可借鉴思路",
                size=30, bold=True, color=COLORS["ink"])
    add_textbox(slide, Inches(0.66), Inches(3.05), Inches(7.7), Inches(0.55),
                "两篇方法思想读书报告：讲原理、来源与可迁移设计", size=16, color=COLORS["muted"])
    add_card(slide, Inches(0.72), Inches(4.22), Inches(5.65), Inches(1.2),
             "Paper A: PNR safety indicator",
             ["Beyond Time to Collision: The Point of No Return as a Reliable Safety Indicator in Rear-End Vehicle Conflicts"],
             COLORS["orange"], COLORS["white"])
    add_card(slide, Inches(6.78), Inches(4.22), Inches(5.65), Inches(1.2),
             "Paper B: Affordance + Attention",
             ["Integrating Affordances and Attention models for Short-Term Object Interaction Anticipation"],
             COLORS["teal"], COLORS["white"])
    add_textbox(slide, Inches(0.68), Inches(6.55), Inches(11.6), Inches(0.28),
                "汇报重点：哪些思想能迁移到纵向/横向 TTC 预测与碰撞概率建模", size=11, color=COLORS["muted"])

    # 2 Storyline
    slide = prs.slides.add_slide(blank)
    slides.append("整体主线")
    add_title(slide, "两篇文章共同指向：TTC 只是时间，概率需要上下文与可避免性", "核心问题")
    add_textbox(slide, Inches(0.70), Inches(1.28), Inches(11.8), Inches(0.5),
                "你的任务不是只估计“多久会接触”，而是判断“是否会发生碰撞”。中间缺的正是风险语义。", size=16, bold=True)
    y = Inches(2.1)
    steps = [
        ("尺度 TTC", "图像尺度变化给出时间接近性", COLORS["blue"]),
        ("PNR", "动作是否还可避免", COLORS["orange"]),
        ("Affordance", "场景中哪里可能冲突", COLORS["teal"]),
        ("Attention", "模型该看谁、看哪里、看哪段历史", COLORS["purple"]),
        ("碰撞概率", "融合时间、空间、可避免性", COLORS["green"]),
    ]
    x = Inches(0.75)
    for i, (t, b, c) in enumerate(steps):
        add_card(slide, x, y, Inches(2.15), Inches(1.35), t, [b], c, COLORS["white"])
        if i < len(steps) - 1:
            arrow(slide, x + Inches(2.17), y + Inches(0.67), x + Inches(2.65), y + Inches(0.67), COLORS["muted"])
        x += Inches(2.38)
    add_textbox(slide, Inches(0.75), Inches(4.15), Inches(11.7), Inches(0.7),
                "组会结论先行：PNR 提供“能不能避开”的物理边界；Affordance-Attention 提供“哪里/谁最可能冲突”的学习框架。",
                size=18, bold=True, color=COLORS["ink"])
    add_source(slide, "整理自两篇用户指定论文主题与已检索 arXiv 信息；PNR 原文未能公开获取，公式为后碰安全建模通用表达")

    # 3 PNR source and question
    slide = prs.slides.add_slide(blank)
    slides.append("PNR来源")
    add_title(slide, "PNR 文章要解决的不是“何时碰”，而是“还来不来得及避开”", "Paper A")
    add_card(slide, Inches(0.70), Inches(1.35), Inches(4.0), Inches(1.55),
             "来源状态", ["题名来自用户提供", "公开检索未获完整 PDF", "以下为概念与公式复原式解析"], COLORS["orange"], COLORS["light"])
    add_card(slide, Inches(5.00), Inches(1.35), Inches(3.5), Inches(1.55),
             "传统 TTC", ["TTC = d / v_rel", "只描述时间接近性"], COLORS["blue"], COLORS["white"])
    add_card(slide, Inches(8.80), Inches(1.35), Inches(3.5), Inches(1.55),
             "PNR", ["Point of No Return", "描述可避免性边界"], COLORS["orange"], COLORS["white"])
    # diagram
    y0 = Inches(4.15)
    base_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.2), y0, Inches(11.0), y0)
    base_line.line.color.rgb = COLORS["muted"]
    base_line.line.width = Pt(1.4)
    for x0, label, c in [(Inches(1.35), "当前位置", COLORS["blue"]), (Inches(5.2), "PNR点", COLORS["orange"]), (Inches(10.2), "碰撞点", COLORS["red"])]:
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x0, y0 - Inches(0.12), Inches(0.24), Inches(0.24))
        set_fill(circ, c)
        circ.line.fill.background()
        add_textbox(slide, x0 - Inches(0.45), y0 + Inches(0.22), Inches(1.3), Inches(0.24), label, size=10, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.3), Inches(4.85), Inches(3.4), Inches(0.35), "还可通过制动/避让避免", size=12, color=COLORS["green"])
    add_textbox(slide, Inches(5.1), Inches(4.85), Inches(4.5), Inches(0.35), "越过后，即使采取最大动作也可能不够", size=12, color=COLORS["red"])
    add_source(slide, "Paper A title: Beyond Time to Collision...；本页示意图为概念改绘")

    # 4 TTC limitation
    slide = prs.slides.add_slide(blank)
    slides.append("TTC局限")
    add_title(slide, "同样 TTC 不等于同样危险：速度尺度被 TTC 抹掉了", "TTC 局限")
    add_textbox(slide, Inches(0.72), Inches(1.33), Inches(5.0), Inches(0.42),
                "TTC 只看比例：", size=15, bold=True)
    add_textbox(slide, Inches(0.72), Inches(1.82), Inches(5.0), Inches(0.6),
                "TTC = d / v_rel", size=26, bold=True, color=COLORS["blue"], font="Aptos")
    add_card(slide, Inches(0.78), Inches(2.80), Inches(4.6), Inches(1.2),
             "场景 A", ["d = 20 m, v_rel = 10 m/s", "TTC = 2 s, a_req = 2.5 m/s²"], COLORS["blue"], COLORS["white"])
    add_card(slide, Inches(0.78), Inches(4.18), Inches(4.6), Inches(1.2),
             "场景 B", ["d = 40 m, v_rel = 20 m/s", "TTC = 2 s, a_req = 5.0 m/s²"], COLORS["red"], COLORS["white"])
    # plot-like visual
    ax = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.25), Inches(1.45), Inches(5.65), Inches(4.55))
    set_fill(ax, COLORS["white"])
    set_line(ax, COLORS["line"])
    add_textbox(slide, Inches(6.45), Inches(1.64), Inches(5.0), Inches(0.28), "PNR 会把速度尺度重新带回来", size=15, bold=True)
    add_textbox(slide, Inches(6.55), Inches(2.25), Inches(4.8), Inches(0.55),
                "a_req = v_rel² / (2d)\n      = v_rel / (2TTC)", size=21, bold=True, color=COLORS["orange"], font="Aptos")
    add_textbox(slide, Inches(6.55), Inches(3.55), Inches(4.8), Inches(0.9),
                "TTC 相同但 v_rel 更大时，所需制动更强；如果超过车辆能力，就进入 PNR。", size=15, color=COLORS["ink"])
    add_textbox(slide, Inches(6.55), Inches(5.10), Inches(4.8), Inches(0.45),
                "这解释了为什么单用 TTC 阈值容易误判风险。", size=15, bold=True, color=COLORS["red"])
    add_source(slide, "公式为后碰制动可避免性通用推导；数值例子为讲解用")

    # 5 PNR formulas
    slide = prs.slides.add_slide(blank)
    slides.append("PNR公式")
    add_title(slide, "PNR 可写成一个连续风险量：所需动作 / 可用动作", "公式原理")
    add_card(slide, Inches(0.70), Inches(1.30), Inches(5.55), Inches(1.3),
             "步骤 1：扣除反应时间", ["d_eff = d - v_rel · t_r - d_0", "d_eff ≤ 0 时，已经没有有效制动距离"], COLORS["orange"], COLORS["white"])
    add_card(slide, Inches(0.70), Inches(2.85), Inches(5.55), Inches(1.3),
             "步骤 2：计算所需减速度", ["a_req = v_rel² / (2 · d_eff)", "它衡量“为了不碰，需要多猛地刹”"], COLORS["orange"], COLORS["white"])
    add_card(slide, Inches(0.70), Inches(4.40), Inches(5.55), Inches(1.3),
             "步骤 3：和车辆能力比较", ["PNR_score = a_req / b_max", "PNR_score > 1 表示越过不可避免边界"], COLORS["red"], COLORS["white"])
    # safe curve
    plot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(1.35), Inches(4.8), Inches(4.55))
    set_fill(plot, COLORS["white"])
    set_line(plot, COLORS["line"])
    add_textbox(slide, Inches(7.22), Inches(1.52), Inches(4.3), Inches(0.28), "状态空间中的 PNR 临界线", size=14, bold=True)
    # axes
    arrow(slide, Inches(7.45), Inches(5.25), Inches(11.35), Inches(5.25), COLORS["ink"], 1.2)
    arrow(slide, Inches(7.45), Inches(5.25), Inches(7.45), Inches(2.0), COLORS["ink"], 1.2)
    add_textbox(slide, Inches(10.25), Inches(5.45), Inches(1.1), Inches(0.22), "v_rel", size=9, font="Aptos")
    add_textbox(slide, Inches(7.0), Inches(2.0), Inches(0.4), Inches(0.22), "d", size=9, font="Aptos")
    points = [(7.6, 4.85), (8.2, 4.55), (8.8, 4.1), (9.4, 3.45), (10.0, 2.65), (10.8, 2.05)]
    for i in range(len(points)-1):
        x1, y1 = points[i]
        x2, y2 = points[i+1]
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)).line.color.rgb = COLORS["orange"]
    add_textbox(slide, Inches(8.35), Inches(2.18), Inches(2.2), Inches(0.25), "未到 PNR", size=11, color=COLORS["green"], align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.35), Inches(4.52), Inches(2.2), Inches(0.25), "已到 PNR", size=11, color=COLORS["red"], align=PP_ALIGN.CENTER)
    add_source(slide, "整理自后碰 PNR/DRAC/制动距离模型；临界线为概念改绘")

    # 6 Scale TTC issue
    slide = prs.slides.add_slide(blank)
    slides.append("尺度TTC与PNR")
    add_title(slide, "只用尺度预测 TTC 时，PNR 不能被严格计算", "与你的任务直接相关")
    add_textbox(slide, Inches(0.74), Inches(1.35), Inches(5.4), Inches(0.55),
                "尺度 TTC 给出：", size=15, bold=True)
    add_textbox(slide, Inches(0.74), Inches(1.90), Inches(5.4), Inches(0.6),
                "TTC = d / v_rel", size=25, bold=True, color=COLORS["blue"], font="Aptos")
    add_textbox(slide, Inches(0.74), Inches(2.78), Inches(5.4), Inches(0.65),
                "PNR 需要：", size=15, bold=True)
    add_textbox(slide, Inches(0.74), Inches(3.32), Inches(5.8), Inches(0.6),
                "a_req = v_rel / (2TTC)\n或 a_req = d / (2TTC²)", size=22, bold=True, color=COLORS["orange"], font="Aptos")
    add_card(slide, Inches(6.65), Inches(1.28), Inches(5.45), Inches(1.2),
             "严格路线", ["补充距离 d 或相对速度 v_rel", "来自深度、雷达、CAN、单目伪深度"], COLORS["green"], COLORS["white"])
    add_card(slide, Inches(6.65), Inches(2.75), Inches(5.45), Inches(1.2),
             "弱物理路线", ["有 ego speed 时用最坏情况", "a_req_worst = v_ego / (2TTC)"], COLORS["orange"], COLORS["white"])
    add_card(slide, Inches(6.65), Inches(4.22), Inches(5.45), Inches(1.2),
             "学习路线", ["不声称精确 PNR", "把 PNR 思想做成 risk feature / loss"], COLORS["purple"], COLORS["white"])
    add_textbox(slide, Inches(0.78), Inches(5.93), Inches(11.3), Inches(0.45),
                "组会建议：你的论文可说“PNR-inspired risk”，不要在无距离/速度时声称精确 PNR。", size=16, bold=True, color=COLORS["red"])
    add_source(slide, "结合用户当前尺度 TTC 设定给出的可实现路径")

    # 7 Paper B task/source
    slide = prs.slides.add_slide(blank)
    slides.append("Affordance来源")
    add_title(slide, "Affordance-Attention 文章预测的是“下一次人-物交互”", "Paper B")
    add_card(slide, Inches(0.70), Inches(1.26), Inches(5.9), Inches(1.6),
             "任务：Short-Term Object Interaction Anticipation",
             ["输入第一视角视频", "预测即将被交互的物体、动作、类别、time-to-contact"], COLORS["teal"], COLORS["light"])
    add_card(slide, Inches(6.95), Inches(1.26), Inches(5.1), Inches(1.6),
             "注意概念差异",
             ["这里的 time-to-contact 是手接触物体", "不是车辆碰撞 TTC"], COLORS["orange"], COLORS["white"])
    # video sequence visual
    for i in range(5):
        x = Inches(1.0 + i*1.05)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.55), Inches(0.82), Inches(0.58))
        set_fill(rect, RGBColor(228, 238, 246))
        set_line(rect, COLORS["blue"], 0.8)
        add_textbox(slide, x + Inches(0.08), Inches(3.74), Inches(0.66), Inches(0.16), f"t-{4-i}", size=8, align=PP_ALIGN.CENTER, font="Aptos")
    arrow(slide, Inches(6.0), Inches(3.84), Inches(7.0), Inches(3.84), COLORS["muted"])
    add_card(slide, Inches(7.15), Inches(3.25), Inches(4.45), Inches(1.25),
             "输出", ["active object + verb/noun + time-to-contact"], COLORS["teal"], COLORS["white"])
    add_textbox(slide, Inches(1.0), Inches(5.2), Inches(10.6), Inches(0.48),
                "核心不是“识别当前发生了什么”，而是提前预测短期未来会和哪个对象发生交互。", size=17, bold=True)
    add_source(slide, "Source: arXiv:2602.14837; arXiv:2406.01194 AFF-ttention")

    # 8 Architecture
    slide = prs.slides.add_slide(blank)
    slides.append("STAformer结构")
    add_title(slide, "STAformer 的关键：当前帧看空间，视频片段看运动", "方法结构")
    # flow
    add_card(slide, Inches(0.75), Inches(1.35), Inches(2.35), Inches(1.1), "当前帧 image", ["清楚定位物体"], COLORS["blue"], COLORS["white"])
    add_card(slide, Inches(0.75), Inches(3.00), Inches(2.35), Inches(1.1), "历史视频 video", ["捕捉手/物体运动"], COLORS["teal"], COLORS["white"])
    add_card(slide, Inches(4.05), Inches(1.60), Inches(3.0), Inches(1.3), "dual image-video attention", ["让空间细节和运动趋势互相选择"], COLORS["purple"], COLORS["light"])
    add_card(slide, Inches(4.05), Inches(3.25), Inches(3.0), Inches(1.3), "frame-guided temporal pooling", ["用当前关键帧指导历史时序聚合"], COLORS["purple"], COLORS["white"])
    add_card(slide, Inches(8.0), Inches(2.15), Inches(3.55), Inches(1.55), "交互预测", ["哪个物体", "什么动作", "多久接触"], COLORS["green"], COLORS["white"])
    arrow(slide, Inches(3.1), Inches(1.9), Inches(4.0), Inches(2.2), COLORS["muted"])
    arrow(slide, Inches(3.1), Inches(3.55), Inches(4.0), Inches(3.85), COLORS["muted"])
    arrow(slide, Inches(7.05), Inches(2.25), Inches(8.0), Inches(2.65), COLORS["muted"])
    arrow(slide, Inches(7.05), Inches(3.78), Inches(8.0), Inches(3.1), COLORS["muted"])
    add_textbox(slide, Inches(0.85), Inches(5.35), Inches(10.8), Inches(0.55),
                "可迁移理解：当前帧决定“目标是谁/在哪里”，历史片段决定“它正在怎样接近”。", size=17, bold=True, color=COLORS["ink"])
    add_source(slide, "Source: arXiv:2602.14837；结构为论文方法概念改绘")

    # 9 Affordance and attention
    slide = prs.slides.add_slide(blank)
    slides.append("Affordance原理")
    add_title(slide, "Affordance 提供“哪里可能发生交互”，Attention 决定“该看哪里”", "核心思想")
    add_card(slide, Inches(0.70), Inches(1.35), Inches(3.45), Inches(1.55),
             "Environment affordance", ["场景里哪些区域常发生交互", "像一张长期交互先验图"], COLORS["teal"], COLORS["light"])
    add_card(slide, Inches(4.93), Inches(1.35), Inches(3.45), Inches(1.55),
             "Interaction hotspot", ["手和物体运动指向哪里", "预测未来接触热点"], COLORS["orange"], COLORS["white"])
    add_card(slide, Inches(9.15), Inches(1.35), Inches(3.05), Inches(1.55),
             "Attention", ["从候选目标中选择关键对象", "让模型少看无关区域"], COLORS["purple"], COLORS["white"])
    # heatmap schematic
    field = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(3.45), Inches(9.3), Inches(2.25))
    set_fill(field, RGBColor(239, 244, 248))
    set_line(field, COLORS["line"])
    for x, y, r, c in [(3.2, 4.2, 0.5, COLORS["teal"]), (6.1, 4.55, 0.72, COLORS["orange"]), (8.8, 4.0, 0.45, COLORS["purple"])]:
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(r), Inches(r))
        set_fill(s, c, 30)
        s.line.fill.background()
    add_textbox(slide, Inches(2.25), Inches(5.92), Inches(8.9), Inches(0.32),
                "风险不是均匀分布的：模型要把概率质量集中到“可交互/可冲突”的区域。", size=15, bold=True, align=PP_ALIGN.CENTER)
    add_source(slide, "Source: arXiv:2602.14837 / 2406.01194；本页为 affordance/hotspot 概念改绘")

    # 10 Transfer to driving
    slide = prs.slides.add_slide(blank)
    slides.append("迁移到驾驶")
    add_title(slide, "把人-物 affordance 改成驾驶场景的 collision affordance", "借鉴点")
    add_card(slide, Inches(0.75), Inches(1.35), Inches(3.4), Inches(1.25),
             "原任务", ["手会和哪个物体交互？"], COLORS["teal"], COLORS["white"])
    add_card(slide, Inches(4.95), Inches(1.35), Inches(3.4), Inches(1.25),
             "你的任务", ["哪辆车/哪个区域会冲突？"], COLORS["blue"], COLORS["white"])
    add_card(slide, Inches(9.15), Inches(1.35), Inches(3.0), Inches(1.25),
             "迁移概念", ["collision affordance map"], COLORS["orange"], COLORS["light"])
    arrow(slide, Inches(4.15), Inches(1.98), Inches(4.95), Inches(1.98), COLORS["muted"])
    arrow(slide, Inches(8.35), Inches(1.98), Inches(9.15), Inches(1.98), COLORS["muted"])
    add_multiline(slide, Inches(1.0), Inches(3.18), Inches(11.2), Inches(1.5),
                  ["同车道前车：纵向碰撞 affordance 高",
                   "邻车道切入：横向碰撞 affordance 高",
                   "路口/遮挡边缘：短时冲突不确定性高"], size=16, bullet=True)
    add_textbox(slide, Inches(1.0), Inches(5.35), Inches(11.2), Inches(0.5),
                "这能补上纯尺度 TTC 的短板：TTC 给时间，affordance 给空间与语义上下文。", size=18, bold=True, color=COLORS["ink"])
    add_source(slide, "本页为跨任务迁移设计：human-object interaction → vehicle conflict anticipation")

    # 11 Proposed framework
    slide = prs.slides.add_slide(blank)
    slides.append("你的框架")
    add_title(slide, "建议框架：尺度 TTC 是主干，PNR/Affordance 是风险增强层", "可落地方案")
    stages = [
        ("视频序列 + 候选车框", "当前帧 crop\n历史尺度变化", COLORS["blue"]),
        ("TTC 分解头", "TTC_long\nTTC_lat", COLORS["teal"]),
        ("风险先验层", "PNR-inspired score\ncollision affordance map", COLORS["orange"]),
        ("Attention 融合", "关注冲突目标\n关注关键时刻", COLORS["purple"]),
        ("概率输出", "P(collision)\n+ hotspot", COLORS["green"]),
    ]
    x = Inches(0.55)
    for i, (t, b, c) in enumerate(stages):
        add_card(slide, x, Inches(1.65), Inches(2.25), Inches(1.45), t, b.split("\n"), c, COLORS["white"])
        if i < len(stages) - 1:
            arrow(slide, x + Inches(2.25), Inches(2.38), x + Inches(2.65), Inches(2.38), COLORS["muted"])
        x += Inches(2.5)
    add_textbox(slide, Inches(0.85), Inches(4.10), Inches(11.5), Inches(0.72),
                "如果没有速度/距离，就不要声称精确 PNR；可以把 PNR 作为启发式 score 或监督约束。", size=17, bold=True, color=COLORS["red"])
    add_multiline(slide, Inches(0.95), Inches(5.05), Inches(11.2), Inches(0.85),
                  ["显式特征：1/TTC_long, 1/TTC_lat, box expansion, overlap, same-lane",
                   "学习目标：TTC 回归 + collision probability + optional hotspot heatmap"], size=14.5, bullet=True)
    add_source(slide, "综合 Paper A 的可避免性思想与 Paper B 的 affordance-attention 结构")

    # 12 What to borrow / avoid
    slide = prs.slides.add_slide(blank)
    slides.append("借鉴与边界")
    add_title(slide, "可借鉴的是建模逻辑，不是照搬任务或公式", "方法边界")
    add_card(slide, Inches(0.72), Inches(1.35), Inches(5.6), Inches(3.85),
             "值得借鉴", ["PNR: 从 TTC 到可避免性", "Affordance: 冲突区域先验", "Attention: 目标/时间片选择", "Fusion: 物理先验 + 学习输出"], COLORS["green"], COLORS["white"])
    add_card(slide, Inches(6.95), Inches(1.35), Inches(5.25), Inches(3.85),
             "不要照搬", ["无速度/距离时不说精确 PNR", "人手轨迹模块不能直接用于车辆", "室内环境记忆不能直接变成道路位置记忆", "不要伪造原文结果或数值"], COLORS["red"], COLORS["light"])
    add_textbox(slide, Inches(1.0), Inches(5.85), Inches(11.0), Inches(0.38),
                "最稳的论文表述：PNR-inspired + traffic affordance-aware collision probability prediction。", size=17, bold=True)
    add_source(slide, "本页为组会讨论结论")

    # 13 Experiments
    slide = prs.slides.add_slide(blank)
    slides.append("实验设计")
    add_title(slide, "实验设计要证明：每个思想都真的改善概率判断", "后续可做")
    rows = [
        ("Baseline", "尺度 TTC_x/TTC_y → 概率", "验证主干能力"),
        ("+ PNR-inspired", "加入 TTC-derived risk / ego speed risk", "看 false-safe 是否下降"),
        ("+ Affordance", "加入 same-lane / overlap / lane prior", "看横向冲突是否更准"),
        ("+ Attention", "时序注意力 + 目标注意力", "看短时预测稳定性"),
        ("Full", "TTC + risk + affordance + attention", "最终碰撞概率"),
    ]
    table = slide.shapes.add_table(len(rows)+1, 3, Inches(0.75), Inches(1.45), Inches(11.85), Inches(3.55)).table
    headers = ["设置", "输入/模块", "要证明什么"]
    for j, htxt in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = htxt
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["ink"]
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = COLORS["white"]
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.name = "Noto Sans CJK SC"
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["white"] if i % 2 else COLORS["light"]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10.5)
                p.font.name = "Noto Sans CJK SC"
                p.font.color.rgb = COLORS["ink"]
    add_multiline(slide, Inches(0.9), Inches(5.35), Inches(11.4), Inches(0.75),
                  ["指标建议：Brier score / ECE 看概率校准，AUROC/F1 看分类，false-safe rate 看安全风险。"], size=15.5, bullet=True)
    add_source(slide, "实验建议面向用户当前 TTC/碰撞概率课题")

    # 14 Summary
    slide = prs.slides.add_slide(blank)
    slides.append("总结")
    add_title(slide, "一句话总结：TTC 预测之后，还需要“可避免性”和“冲突上下文”", "Take-home")
    add_textbox(slide, Inches(0.85), Inches(1.35), Inches(11.6), Inches(0.78),
                "PNR 回答“是否还避得开”；Affordance-Attention 回答“哪里/谁最可能发生交互”。", size=23, bold=True, color=COLORS["ink"])
    add_card(slide, Inches(1.0), Inches(2.70), Inches(3.3), Inches(1.35),
             "Paper A", ["TTC → 可避免性边界", "适合做 risk score"], COLORS["orange"], COLORS["white"])
    add_card(slide, Inches(5.0), Inches(2.70), Inches(3.3), Inches(1.35),
             "Paper B", ["Affordance + Attention", "适合做冲突区域建模"], COLORS["teal"], COLORS["white"])
    add_card(slide, Inches(9.0), Inches(2.70), Inches(3.3), Inches(1.35),
             "你的方向", ["纵/横向 TTC", "碰撞概率 + hotspot"], COLORS["green"], COLORS["white"])
    add_textbox(slide, Inches(1.0), Inches(5.42), Inches(11.3), Inches(0.48),
                "推荐主张：从 scale-based TTC estimation 走向 risk-aware collision anticipation。", size=18, bold=True, color=COLORS["blue"], align=PP_ALIGN.CENTER)
    add_source(slide, "Sources: user-provided PNR paper title; arXiv:2602.14837; arXiv:2406.01194")

    prs.save(PPTX_PATH)
    return slides


def write_supporting_files(slides):
    outline = """# 组会 PPT 大纲

检测 paper_type: methods。

中心论点：两篇文章都可用于把尺度 TTC 预测升级为 risk-aware collision anticipation。

术语表：

| Canonical term | First-use definition | Decision |
|---|---|---|
| TTC | Time to Collision / Time-to-Contact，按上下文区分 | 车辆场景用 TTC，Paper B 用 time-to-contact |
| PNR | Point of No Return | 保留英文缩写，中文解释为不可避免边界 |
| Affordance | 可供性/可交互性先验 | PPT 中保留英文，中文解释 |
| Collision affordance map | 驾驶场景冲突可供性图 | 作为迁移概念使用 |
| STA | Short-Term Object Interaction Anticipation | Paper B 任务名 |
| STAformer | Paper B 的时空注意力模型 | 保留原名 |

幻灯片：
""" + "\n".join([f"{i+1}. {name}" for i, name in enumerate(slides)])
    OUTLINE_PATH.write_text(outline, encoding="utf-8")

    notes = """# 组会讲稿提示

1. 封面：说明这不是普通文献复述，而是把两篇文章转化成你课题可借鉴的方法设计。
2. 整体主线：先强调 TTC 只是时间指标，碰撞概率需要更多风险语义。
3. PNR 来源：说明 PNR 原文未公开获取，本文用其题名主张和通用后碰模型解释原理。
4. TTC 局限：用两个 TTC 都是 2 秒但所需减速度不同的例子，让听众直观看到 TTC 的盲点。
5. PNR 公式：重点讲 d_eff、a_req、PNR_score 三步；PNR_score > 1 是可解释风险边界。
6. 尺度 TTC 与 PNR：这是和你最相关的一页；只靠尺度 TTC 不能严格算 PNR，需要距离或速度。
7. Affordance 来源：提醒 Paper B 的 time-to-contact 是手接触物体，不是车辆 TTC。
8. STAformer 结构：讲 current image 负责空间，video clip 负责运动。
9. Affordance 原理：environment affordance 是长期先验，hotspot 是短期交互焦点。
10. 迁移到驾驶：把厨房里的“哪里可交互”变成道路里的“哪里可能冲突”。
11. 你的框架：这是建议方案页，可以作为后续开题/论文方法雏形。
12. 借鉴与边界：主动说明哪些不能照搬，避免组会上被质疑。
13. 实验设计：强调要用消融证明每个模块的作用，尤其是概率校准和 false-safe。
14. 总结：用一句话收束：TTC + 可避免性 + 冲突上下文。
"""
    NOTES_PATH.write_text(notes, encoding="utf-8")


def audit_pptx():
    prs = Presentation(PPTX_PATH)
    issues = []
    slide_w, slide_h = prs.slide_width, prs.slide_height
    text_counts = []
    for i, slide in enumerate(prs.slides, start=1):
        chars = 0
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > slide_w or shape.top + shape.height > slide_h:
                issues.append(("medium", i, "shape out of bounds"))
            if getattr(shape, "has_text_frame", False):
                txt = shape.text.strip()
                chars += len(txt)
        text_counts.append((i, chars))
        if chars > 520:
            issues.append(("medium", i, f"text-heavy slide: {chars} chars"))
    media_count = 0
    with zipfile.ZipFile(PPTX_PATH) as zf:
        media_count = len([n for n in zf.namelist() if n.startswith("ppt/media/")])
    return len(prs.slides), media_count, issues, text_counts


def write_qa(slide_count, media_count, issues):
    corrected = [
        "全部页面采用可编辑 PPT native 图形，避免低清截图。",
        "对 PNR 原文不可访问状态做了显式标注，未伪造实验结果或原图。",
        "每页控制为一个主观点，并把公式/流程拆成多页讲解。",
        "已重新打开 PPTX 并检查 slide count、media count、shape bounds、文本密度。"
    ]
    issue_lines = "\n".join([f"- {sev}: slide {idx} — {msg}" for sev, idx, msg in issues]) or "- 未发现 high/medium 程序化问题。"
    qa = f"""# QA Report

## 输出状态
- PPTX: `{PPTX_PATH}`
- Slide count: {slide_count}
- Embedded media count: {media_count}
- Supporting notes: `{NOTES_PATH}`
- Outline: `{OUTLINE_PATH}`

## 来源与限制
- Paper A: 用户提供题名 `Beyond Time to Collision: The Point of No Return as a Reliable Safety Indicator in Rear-End Vehicle Conflicts`。公开检索未获得完整原文/PDF，因此 PPT 不使用原文图和实验数值；PNR 公式页基于后碰可避免性、DRAC、制动距离的通用物理模型整理。
- Paper B: `Integrating Affordances and Attention models for Short-Term Object Interaction Anticipation`，arXiv:2602.14837；同时引用 2024 AFF-ttention arXiv:2406.01194 作为来源背景。

## 自检发现
{issue_lines}

## 已修正/设计处理
""" + "\n".join([f"- {x}" for x in corrected]) + """

## 已知限制
- 未做整页渲染预览；当前验证基于 python-pptx 重新打开、结构检查和文本密度检查。
- 若后续拿到 Paper A 原文 PDF，可将概念改绘页替换为原文 Figure，并补充真实结果页。
"""
    QA_PATH.write_text(qa, encoding="utf-8")


if __name__ == "__main__":
    slides = make_deck()
    write_supporting_files(slides)
    slide_count, media_count, issues, _ = audit_pptx()
    write_qa(slide_count, media_count, issues)
    print(f"wrote {PPTX_PATH} with {slide_count} slides")
    print(f"media_count={media_count}")
    print(f"issues={len(issues)}")
